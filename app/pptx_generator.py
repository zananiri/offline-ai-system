"""
Generates PowerPoint (.pptx) files on demand from a plain-language prompt.

Uses the same local Ollama model (qwen2.5) already used elsewhere in this
app for non-translation tasks, plus python-pptx (MIT licensed, fully
offline -- no network calls, no paid API) to actually build the file.

Flow:
    prompt (topic, optionally with attached-document context / slide count)
        -> qwen2.5 asked for a strict JSON outline (title + per-slide
           heading/bullets)
        -> python-pptx turns that outline into a real .pptx file on disk
"""
import json
import re
from pathlib import Path

import ollama
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Kept in sync with main.py's default chat model -- this is the model used
# everywhere else in the app for non-translation work (summarizing,
# classifying invoices, etc.), so presentation generation reuses it too.
OLLAMA_MODEL = "qwen2.5:7b-instruct-q4_K_M"

_OUTLINE_SYSTEM_PROMPT = (
    "You create outlines for PowerPoint presentations. Given a topic (and, "
    "sometimes, source material to summarize), respond with ONLY valid JSON, "
    "no other text, no markdown code fences, in exactly this shape:\n"
    '{"title": string, "subtitle": string, "slides": '
    '[{"heading": string, "bullets": [string, ...]}, ...]}\n\n'
    "\"title\" is a short, presentation-worthy title (not a restatement of "
    "the raw request). \"subtitle\" is a one-line subtitle or tagline (can "
    "be an empty string). Each slide's \"bullets\" list should have 3-5 "
    "short bullet points (each under ~15 words) -- never full paragraphs. "
    "Produce between 5 and 10 content slides unless the user specifies a "
    "different number, in which case follow that instead. If source "
    "material is provided, base the slides on it; otherwise use your own "
    "knowledge of the topic."
)

# Unicode ranges for Hebrew/Arabic script -- same heuristic app/document.py
# uses elsewhere in this project to decide text direction.
_RTL_CHAR_RE = re.compile(r"[\u0590-\u06ff]")


def _is_rtl_text(text: str, threshold: float = 0.3) -> bool:
    letters = [ch for ch in text if ch.isalpha()]
    if not letters:
        return False
    rtl_letters = sum(1 for ch in letters if _RTL_CHAR_RE.match(ch))
    return (rtl_letters / len(letters)) >= threshold


def _strip_json_fences(text: str) -> str:
    text = text.strip()
    if text.startswith("```"):
        text = re.sub(r"^```[a-zA-Z]*\n?", "", text)
        text = re.sub(r"```\s*$", "", text)
    return text.strip()


def _normalize_outline(raw_outline: dict) -> dict | None:
    if not isinstance(raw_outline, dict) or "slides" not in raw_outline:
        return None
    slides = []
    for s in raw_outline.get("slides", []) or []:
        if not isinstance(s, dict):
            continue
        heading = str(s.get("heading", "")).strip()
        bullets = [str(b).strip() for b in (s.get("bullets") or []) if str(b).strip()]
        if heading or bullets:
            slides.append({"heading": heading or "Untitled", "bullets": bullets})
    if not slides:
        return None
    return {
        "title": str(raw_outline.get("title") or "Untitled Presentation"),
        "subtitle": str(raw_outline.get("subtitle") or ""),
        "slides": slides,
    }


def generate_outline(prompt: str, model: str = OLLAMA_MODEL) -> dict:
    """Asks the LLM for a structured slide outline and returns it as a dict.
    Raises ValueError if the model's response can't be parsed into the
    expected shape even after one retry."""
    messages = [
        {"role": "system", "content": _OUTLINE_SYSTEM_PROMPT},
        {"role": "user", "content": prompt},
    ]

    for attempt in range(2):
        response = ollama.chat(model=model, format="json", messages=messages)
        raw = response["message"]["content"]
        try:
            parsed = json.loads(_strip_json_fences(raw))
        except json.JSONDecodeError:
            parsed = None
        outline = _normalize_outline(parsed) if parsed is not None else None
        if outline:
            return outline
        messages.append({
            "role": "user",
            "content": "That wasn't valid. Respond again with ONLY the JSON object, nothing else.",
        })

    raise ValueError("Could not parse a slide outline from the model's response.")


# --- Simple, clean built-in look (no external template file needed) ---
_TITLE_COLOR = RGBColor(0x1A, 0x1A, 0x1A)
_ACCENT_COLOR = RGBColor(0xF9, 0x73, 0x16)  # matches this project's orange accent
_BODY_COLOR = RGBColor(0x33, 0x33, 0x33)


def _style_paragraph_runs(paragraph, size_pt, color, rtl):
    paragraph.alignment = PP_ALIGN.RIGHT if rtl else PP_ALIGN.LEFT
    for run in paragraph.runs:
        run.font.size = Pt(size_pt)
        run.font.color.rgb = color


def build_pptx(outline: dict, output_path: str) -> str:
    """Builds a .pptx file from an outline dict (see generate_outline's
    return shape) using python-pptx."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_text = outline.get("title", "Untitled Presentation")
    subtitle_text = outline.get("subtitle", "")
    deck_rtl = _is_rtl_text(f"{title_text} {subtitle_text}")

    # --- Title slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    _style_paragraph_runs(slide.shapes.title.text_frame.paragraphs[0], 40, _ACCENT_COLOR, deck_rtl)
    for run in slide.shapes.title.text_frame.paragraphs[0].runs:
        run.font.bold = True

    if len(slide.placeholders) > 1 and subtitle_text:
        subtitle_ph = slide.placeholders[1]
        subtitle_ph.text = subtitle_text
        _style_paragraph_runs(subtitle_ph.text_frame.paragraphs[0], 20, _BODY_COLOR, deck_rtl)

    # --- Content slides ---
    content_layout = prs.slide_layouts[1]  # "Title and Content"
    for slide_data in outline.get("slides", []):
        slide = prs.slides.add_slide(content_layout)
        heading = slide_data.get("heading", "")
        bullets = slide_data.get("bullets", [])
        slide_rtl = _is_rtl_text(f"{heading} {' '.join(bullets)}")

        slide.shapes.title.text = heading
        _style_paragraph_runs(slide.shapes.title.text_frame.paragraphs[0], 30, _TITLE_COLOR, slide_rtl)
        for run in slide.shapes.title.text_frame.paragraphs[0].runs:
            run.font.bold = True

        body_ph = slide.placeholders[1]
        tf = body_ph.text_frame
        tf.word_wrap = True
        tf.clear()
        for i, bullet in enumerate(bullets):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = bullet
            _style_paragraph_runs(para, 20, _BODY_COLOR, slide_rtl)

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path


def generate_pptx(prompt: str, output_path: str, model: str = OLLAMA_MODEL) -> str:
    """Full pipeline: prompt -> LLM outline -> .pptx file on disk. Returns
    output_path for convenience."""
    outline = generate_outline(prompt, model=model)
    return build_pptx(outline, output_path)